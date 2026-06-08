from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = PROJECT_ROOT / "CV Final Project Slide Deck.pptx"
OUT_PATH = PROJECT_ROOT / "Problem_A_sample_slide_deck_v2.pptx"
SLIDE_PACKAGE = PROJECT_ROOT / "problem_a_slide_package"


PAPER = RGBColor(251, 250, 246)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(29, 39, 51)
INK = RGBColor(42, 52, 63)
MUTED = RGBColor(91, 103, 117)
LINE = RGBColor(218, 224, 230)
BLUE = RGBColor(64, 107, 168)
GREEN = RGBColor(77, 138, 98)
AMBER = RGBColor(195, 131, 47)
RED = RGBColor(189, 90, 82)
PURPLE = RGBColor(139, 97, 183)
TEAL = RGBColor(57, 138, 134)
GRAY = RGBColor(102, 113, 125)


def clear_slides(prs: Presentation) -> None:
    """Keep template theme/layouts, remove template example slides."""
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def emu_to_in(value: int) -> float:
    return value / 914400


def set_bg(slide, color=PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text: str, x, y, w, h, size=11, bold=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def title(slide, text: str, subtitle: str = "", section: str = "Problem A") -> None:
    textbox(slide, text, Inches(0.38), Inches(0.24), Inches(7.7), Inches(0.34), 19, True, NAVY)
    if subtitle:
        textbox(slide, subtitle, Inches(0.4), Inches(0.61), Inches(8.3), Inches(0.28), 8.8, False, MUTED)
    if section:
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(0.28), Inches(0.95), Inches(0.25))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(235, 239, 244)
        chip.line.color.rgb = RGBColor(210, 217, 225)
        tf = chip.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = section
        r.font.name = "Aptos"
        r.font.size = Pt(7)
        r.font.bold = True
        r.font.color.rgb = GRAY


def footer(slide, page: int, label="Dog Replacement Project") -> None:
    textbox(slide, label, Inches(0.4), Inches(5.38), Inches(4.2), Inches(0.15), 6.5, False, MUTED)
    textbox(slide, f"{page:02d}", Inches(9.15), Inches(5.38), Inches(0.45), Inches(0.15), 6.5, True, MUTED, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, head: str, body: str, accent=BLUE, head_size=10, body_size=8.2):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(0.8)
    textbox(slide, head, x + Inches(0.12), y + Inches(0.10), w - Inches(0.24), Inches(0.18), head_size, True, accent)
    textbox(slide, body, x + Inches(0.12), y + Inches(0.33), w - Inches(0.24), h - Inches(0.38), body_size, False, INK)
    return box


def bullets(slide, items: list[str], x, y, w, h, size=10.5) -> None:
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(3)


def picture(slide, path: Path, x, y, w, h, caption: str | None = None, border=True):
    if not path.exists():
        card(slide, x, y, w, h, "Missing", str(path.relative_to(PROJECT_ROOT)), RED)
        return None
    with Image.open(path) as img:
        iw, ih = img.size
    aspect = iw / ih
    box_aspect = w / h
    if aspect > box_aspect:
        pw = w
        ph = w / aspect
    else:
        ph = h
        pw = h * aspect
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    if border:
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = LINE
        bg.line.width = Pt(0.8)
    pic = slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)
    if caption:
        textbox(slide, caption, x, y + h + Inches(0.02), w, Inches(0.15), 6.5, False, MUTED, PP_ALIGN.CENTER)
    return pic


def arrow(slide, x, y, w=0.32) -> None:
    textbox(slide, "→", x, y, Inches(w), Inches(0.22), 15, True, GRAY, PP_ALIGN.CENTER)


def add_pipeline_cards(slide) -> None:
    items = [
        ("Inputs", "original + replacement", BLUE),
        ("Pose", "24 dog keypoints", GREEN),
        ("Segment", "YOLO/SAM masks", AMBER),
        ("Gate", "feasibility", PURPLE),
        ("Warp", "scale + pose", RED),
        ("Diffuse", "local repair", TEAL),
        ("Score", "cycle metrics", GRAY),
    ]
    x0, y0 = Inches(0.45), Inches(1.1)
    w, h, gap = Inches(1.12), Inches(0.62), Inches(0.18)
    for i, (head, body, color) in enumerate(items):
        x = x0 + i * (w + gap)
        card(slide, x, y0, w, h, head, body, color, 8.2, 6.4)
        if i < len(items) - 1:
            arrow(slide, x + w, y0 + Inches(0.19), 0.17)


def add_cycle_diagram(slide) -> None:
    card(slide, Inches(0.55), Inches(1.2), Inches(1.55), Inches(0.65), "A", "original scene", BLUE)
    card(slide, Inches(0.55), Inches(2.45), Inches(1.55), Inches(0.65), "B", "replacement dog", GREEN)
    card(slide, Inches(3.0), Inches(1.85), Inches(2.0), Inches(0.72), "A + B → A_B", "forward replacement", AMBER)
    card(slide, Inches(5.8), Inches(1.85), Inches(2.15), Inches(0.72), "A_reconstructed", "replace original dog back", RED)
    card(slide, Inches(3.85), Inches(3.45), Inches(3.05), Inches(0.75), "Compare to A", "boundary, edited, dog, background metrics", GRAY)
    textbox(slide, "↘", Inches(2.15), Inches(1.65), Inches(0.45), Inches(0.3), 18, True, GRAY, PP_ALIGN.CENTER)
    textbox(slide, "↗", Inches(2.15), Inches(2.38), Inches(0.45), Inches(0.3), 18, True, GRAY, PP_ALIGN.CENTER)
    textbox(slide, "→", Inches(5.1), Inches(2.05), Inches(0.42), Inches(0.25), 15, True, GRAY, PP_ALIGN.CENTER)
    textbox(slide, "↓", Inches(6.65), Inches(2.72), Inches(0.35), Inches(0.45), 15, True, GRAY, PP_ALIGN.CENTER)


def add_deployment(slide) -> None:
    xcols = [Inches(0.55), Inches(3.75), Inches(6.95)]
    labels = ["Input", "GPU worker", "Artifacts + ops"]
    for x, lab in zip(xcols, labels):
        textbox(slide, lab, x, Inches(1.05), Inches(2.2), Inches(0.18), 8.5, True, MUTED)
    card(slide, xcols[0], Inches(1.35), Inches(2.45), Inches(0.58), "Upload pair", "original + replacement dog", BLUE, 8.5, 6.4)
    card(slide, xcols[0], Inches(2.2), Inches(2.45), Inches(0.58), "Candidate search", "pose-indexed replacement pool", GREEN, 8.5, 6.4)
    card(slide, xcols[1], Inches(1.35), Inches(2.45), Inches(0.58), "Preflight gate", "pose, coverage, scale fit", AMBER, 8.5, 6.4)
    card(slide, xcols[1], Inches(2.2), Inches(2.45), Inches(0.58), "Pipeline", "segment → warp → diffuse", PURPLE, 8.5, 6.4)
    card(slide, xcols[1], Inches(3.05), Inches(2.45), Inches(0.58), "Model cache", "YOLO, SAM, GAN, diffusion", RED, 8.5, 6.4)
    card(slide, xcols[2], Inches(1.35), Inches(2.45), Inches(0.58), "Outputs", "finals, masks, metadata", TEAL, 8.5, 6.4)
    card(slide, xcols[2], Inches(2.2), Inches(2.45), Inches(0.58), "Evaluation", "cycle metrics + visuals", GRAY, 8.5, 6.4)
    card(slide, xcols[2], Inches(3.05), Inches(2.45), Inches(0.58), "Maintenance", "thresholds, failures, retrain", AMBER, 8.5, 6.4)
    textbox(slide, "→", Inches(3.05), Inches(2.0), Inches(0.35), Inches(0.25), 14, True, GRAY, PP_ALIGN.CENTER)
    textbox(slide, "→", Inches(6.25), Inches(2.0), Inches(0.35), Inches(0.25), 14, True, GRAY, PP_ALIGN.CENTER)
    textbox(slide, "↺", Inches(6.35), Inches(3.7), Inches(0.35), Inches(0.25), 14, True, GRAY, PP_ALIGN.CENTER)


def make_deck() -> None:
    prs = Presentation(str(TEMPLATE_PATH))
    clear_slides(prs)
    blank = prs.slide_layouts[10]
    print(f"Template size: {emu_to_in(prs.slide_width):.2f} x {emu_to_in(prs.slide_height):.2f} in")
    p = 1

    # 1. Problem.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Pose-Aware Dog Replacement", "Replace a dog while preserving pose plausibility and local scene consistency.", "")
    picture(slide, SLIDE_PACKAGE / "page_01_problem/01_original_scene.png", Inches(0.45), Inches(1.05), Inches(2.65), Inches(2.75), "Original")
    textbox(slide, "+", Inches(3.15), Inches(2.25), Inches(0.25), Inches(0.22), 17, True, AMBER, PP_ALIGN.CENTER)
    picture(slide, SLIDE_PACKAGE / "page_01_problem/02_replacement_dog.png", Inches(3.45), Inches(1.05), Inches(2.65), Inches(2.75), "Replacement")
    textbox(slide, "→", Inches(6.18), Inches(2.25), Inches(0.25), Inches(0.22), 17, True, AMBER, PP_ALIGN.CENTER)
    picture(slide, SLIDE_PACKAGE / "page_01_problem/03_final_replacement.png", Inches(6.55), Inches(1.05), Inches(2.65), Inches(2.75), "Final")
    card(slide, Inches(0.75), Inches(4.2), Inches(8.45), Inches(0.62), "Core framing", "This is a staged local image-editing system: pose-index candidates, segment, warp conservatively, inpaint locally, then evaluate with cycle consistency.", BLUE)
    footer(slide, p)
    p += 1

    # 2. Pipeline.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Full Replacement Pipeline", "A modular pipeline makes failures inspectable and model variants comparable.")
    add_pipeline_cards(slide)
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/03_original_removal_mask.png", Inches(0.55), Inches(2.15), Inches(1.75), Inches(1.75), "remove dog")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/06_pose_after_warp.png", Inches(2.55), Inches(2.15), Inches(2.25), Inches(1.75), "align pose")
    picture(slide, SLIDE_PACKAGE / "page_08_diffusion_results/02_commit_mask.png", Inches(5.05), Inches(2.15), Inches(1.75), Inches(1.75), "commit mask")
    picture(slide, SLIDE_PACKAGE / "page_08_diffusion_results/06_final_after_postprocess.png", Inches(7.05), Inches(2.15), Inches(2.25), Inches(1.75), "final")
    bullets(slide, ["Candidate quality first", "Geometry correction is safety-limited", "Diffusion repairs seams, not the entire image"], Inches(0.8), Inches(4.35), Inches(8.4), Inches(0.5), 9.5)
    footer(slide, p)
    p += 1

    # 3. Evaluation.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Evaluation First: Cycle Replacement", "No perfect target image exists, so we reconstruct the original through a forward-back cycle.")
    add_cycle_diagram(slide)
    picture(slide, SLIDE_PACKAGE / "page_04_composite_score/01_best_variant_summary.png", Inches(0.65), Inches(4.35), Inches(3.85), Inches(0.75), "best variant by group")
    picture(slide, SLIDE_PACKAGE / "page_04_composite_score/04_quality_runtime_tradeoff.png", Inches(5.0), Inches(3.55), Inches(4.25), Inches(1.55), "quality/runtime tradeoff")
    bullets(slide, ["Metrics focus on edited, dog, seam, and nearby background regions.", "Runtime is shown separately, not hidden inside the score."], Inches(8.15), Inches(1.35), Inches(1.25), Inches(2.35), 8.5)
    footer(slide, p)
    p += 1

    # 4. Dog-pose training.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Dog-Pose Model Training And Filtering", "The custom pose model gates bad replacements before expensive editing.")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_pose_training/02_dog_pose_training_setup_card.png", Inches(0.5), Inches(1.05), Inches(4.3), Inches(1.65), "training setup")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_pose_training/01_dog_pose_training_curves_available_log.png", Inches(5.05), Inches(1.05), Inches(4.35), Inches(1.65), "local training log")
    picture(slide, SLIDE_PACKAGE / "page_05_dog_pose_training/03_pose_filter_distribution.png", Inches(0.5), Inches(3.05), Inches(4.3), Inches(1.65), "filter distribution")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_pose_validation/01_dog_pose_validation_examples.png", Inches(5.05), Inches(3.05), Inches(4.35), Inches(1.65), "validation examples")
    footer(slide, p)
    p += 1

    # 5. Main branch.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Main Branch: Segmentation + Conservative Warp", "The main branch avoids extreme deformation and lets diffusion repair only local artifacts.")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/01_original_image.png", Inches(0.45), Inches(1.05), Inches(1.7), Inches(1.35), "original")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/04_replacement_conservative_mask.png", Inches(2.35), Inches(1.05), Inches(1.7), Inches(1.35), "mask")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/05_pose_before_warp.png", Inches(4.25), Inches(1.05), Inches(2.3), Inches(1.35), "before")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/06_pose_after_warp.png", Inches(6.85), Inches(1.05), Inches(2.3), Inches(1.35), "after")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/07_replacement_before_after_pose.png", Inches(0.75), Inches(3.0), Inches(4.0), Inches(1.75), "replacement skeleton before/after")
    picture(slide, SLIDE_PACKAGE / "page_06_main_pipeline/08_seed_composite_before_diffusion.png", Inches(5.25), Inches(3.0), Inches(4.0), Inches(1.75), "seed before diffusion")
    footer(slide, p)
    p += 1

    # 6. Face DragGAN.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Face-DragGAN Branch: Inversion + PTI", "A slower challenger branch for semantically meaningful dog-face adjustment.")
    face_imgs = [
        ("crop", "01_replacement_face_crop.png"),
        ("inversion", "02_stylegan_inversion_projection.png"),
        ("PTI", "03_pti_projection.png"),
        ("handles", "04_draggan_handle_plan.png"),
        ("edited", "06_draggan_edited_face.png"),
    ]
    for i, (cap, name) in enumerate(face_imgs):
        picture(slide, SLIDE_PACKAGE / f"page_07_face_draggan/{name}", Inches(0.55 + i * 1.83), Inches(1.05), Inches(1.55), Inches(1.45), cap)
        if i < len(face_imgs) - 1:
            arrow(slide, Inches(2.05 + i * 1.83), Inches(1.67), 0.22)
    card(slide, Inches(0.7), Inches(3.0), Inches(3.55), Inches(0.85), "Why PTI?", "Inversion maps an arbitrary dog face into AFHQ-dog StyleGAN space; PTI fine-tunes around that face to preserve identity before DragGAN editing.", PURPLE, 9.5, 7.2)
    picture(slide, SLIDE_PACKAGE / "page_07_face_draggan/07_face_draggan_seed_composite.png", Inches(4.6), Inches(2.85), Inches(2.15), Inches(1.65), "reinserted seed")
    picture(slide, SLIDE_PACKAGE / "page_07_face_draggan/08_face_draggan_final.png", Inches(7.05), Inches(2.85), Inches(2.15), Inches(1.65), "final")
    footer(slide, p)
    p += 1

    # 7. Diffusion.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Local Diffusion And Harmonization", "Generation can use context, but only a smaller commit region is blended back.")
    imgs = [
        ("generation", "01_generation_mask.png"),
        ("commit", "02_commit_mask.png"),
        ("seed", "03_seed_before_diffusion.png"),
        ("diffusion", "04_realvisxl_refined.png"),
        ("final", "06_final_after_postprocess.png"),
    ]
    for i, (cap, name) in enumerate(imgs):
        picture(slide, SLIDE_PACKAGE / f"page_08_diffusion_results/{name}", Inches(0.5 + i * 1.82), Inches(1.0), Inches(1.5), Inches(1.25), cap)
    card(slide, Inches(0.65), Inches(2.8), Inches(3.0), Inches(0.75), "Shown model", "realvisxl_inpaint, 28 steps, max side 1024, guidance 6.5, strength 0.72; dog-body generation disabled.", TEAL, 9, 7)
    picture(slide, SLIDE_PACKAGE / "page_08_diffusion_results/07_diffusion_composite_score.png", Inches(4.0), Inches(2.65), Inches(2.55), Inches(1.75), "diffusion quality")
    picture(slide, SLIDE_PACKAGE / "page_08_diffusion_results/08_diffusion_runtime.png", Inches(6.95), Inches(2.65), Inches(2.55), Inches(1.75), "runtime")
    footer(slide, p)
    p += 1

    # 8. Results summary.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Comparison Summary", "Cycle-consistency gives a repeatable pilot benchmark for design choices.")
    picture(slide, SLIDE_PACKAGE / "appendix/comparison_charts/02_group_a_composite_score.png", Inches(0.45), Inches(1.05), Inches(2.85), Inches(1.65), "segmentation")
    picture(slide, SLIDE_PACKAGE / "appendix/comparison_charts/02_group_b_composite_score.png", Inches(3.55), Inches(1.05), Inches(2.85), Inches(1.65), "pose")
    picture(slide, SLIDE_PACKAGE / "appendix/comparison_charts/02_group_e_composite_score.png", Inches(6.65), Inches(1.05), Inches(2.85), Inches(1.65), "Face-DragGAN")
    bullets(slide, [
        "Segmentation and mask choice strongly affect seam quality.",
        "Pose adjustment helps only when deformation stays conservative.",
        "Face-DragGAN adds semantic control but costs runtime and may drift.",
        "Diffusion helps most when masks avoid rewriting the dog body.",
    ], Inches(0.85), Inches(3.35), Inches(8.1), Inches(1.0), 10.5)
    footer(slide, p)
    p += 1

    # 9. Backup pair 1.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Backup Example: Main Pipeline Pair 1", "Extra visual result if the audience asks for another case.", "Appendix")
    picture(slide, SLIDE_PACKAGE / "appendix/backup_examples/pair_1_main_pipeline/00_backup_overview_contact_sheet.png", Inches(0.45), Inches(1.0), Inches(9.1), Inches(4.15), border=False)
    footer(slide, p, "Backup visual")
    p += 1

    # 10. Backup pair 3.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Backup Example: Face-DragGAN Pair", "Extra visual result for the face-specific branch.", "Appendix")
    picture(slide, SLIDE_PACKAGE / "appendix/backup_examples/pair_3_face_draggan_pipeline/00_backup_overview_contact_sheet.png", Inches(0.45), Inches(1.0), Inches(9.1), Inches(4.15), border=False)
    footer(slide, p, "Backup visual")
    p += 1

    # 11. EDA.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Appendix: Dataset / Pose-Index EDA", "Counts, distributions, geometry signals, and sample outliers.", "Appendix")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_dataset_eda/01_record_counts_and_schema_summary.png", Inches(0.45), Inches(1.0), Inches(2.1), Inches(1.65), "records")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_dataset_eda/02_dataset_label_distributions.png", Inches(2.8), Inches(1.0), Inches(3.0), Inches(1.65), "labels")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_dataset_eda/03_image_geometry_and_pose_quality.png", Inches(6.05), Inches(1.0), Inches(3.0), Inches(1.65), "geometry")
    picture(slide, SLIDE_PACKAGE / "appendix/generated_dataset_eda/04_raw_samples_and_outliers.png", Inches(0.85), Inches(3.1), Inches(8.3), Inches(1.65), "samples and outliers")
    footer(slide, p, "Appendix")
    p += 1

    # 12. Deployment.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Appendix: Model Operations Plan", "A GPU batch deployment path for the current notebook workflow.", "Appendix")
    add_deployment(slide)
    bullets(slide, ["Maintenance loop: review failure cases, tune thresholds, rerun cycle regression tests, retrain dog-pose model when failures cluster."], Inches(0.8), Inches(4.45), Inches(8.4), Inches(0.35), 9)
    footer(slide, p, "Appendix")
    p += 1

    # 13. Model reference.
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "Appendix: Model Families Used", "Compact reference map for Q&A.", "Appendix")
    models = [
        ("YOLO", "dog detect / seg / custom pose", BLUE),
        ("SAM", "prompted mask challenger", GREEN),
        ("StyleGAN2", "AFHQ dog-face latent space", PURPLE),
        ("PTI", "personalized generator tuning", AMBER),
        ("DragGAN", "nose-guided latent editing", RED),
        ("Diffusion", "local inpainting repair", TEAL),
    ]
    for i, (head, body, color) in enumerate(models):
        row, col = divmod(i, 3)
        card(slide, Inches(0.75 + col * 3.0), Inches(1.25 + row * 1.05), Inches(2.35), Inches(0.72), head, body, color, 9.5, 6.8)
    bullets(slide, [
        "Reference links are in problem_a_slide_package/appendix/model_references/README.md.",
        "The deck uses self-generated schematics rather than copied paper figures.",
    ], Inches(1.0), Inches(4.05), Inches(8.0), Inches(0.55), 9.5)
    footer(slide, p, "Appendix")

    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    make_deck()
